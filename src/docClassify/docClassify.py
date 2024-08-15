from src.spark.SparkApi import SparkLLM
import time
import json5
import os
from docx import Document
from pptx import Presentation
import olefile
import re

with open('../config.json', encoding='utf-8') as f:
    config = json5.load(f)

appid = 'e76d7d8f'
api_secret = 'Y2Y2ODc2OGQyOWFjMWZhY2JkOTllMDVl'
api_key = '990e2770b030441fbcc126c691daf5cd'

domain = "generalv3.5"    # v3.0版本
Spark_url = "wss://spark-api.xf-yun.com/v3.5/chat"  # v3.5环服务地址


def docClassify(file_path, max_length, keyword_list):
    try:
        file_extension = os.path.splitext(file_path)[1].lower()

        if file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read(max_length)

        elif file_extension == '.docx':
            doc = Document(file_path)
            content = ""
            for para in doc.paragraphs:
                content += para.text + '\n'
                if len(content) >= max_length:
                    break
            content = content[:max_length]

        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + '\n'
                        if len(content) >= max_length:
                            break
                if len(content) >= max_length:
                    break
            content = content[:max_length]

        elif file_extension == '.doc':
            if olefile.isOleFile(file_path):
                ole = olefile.OleFileIO(file_path)
                with ole.openstream('WordDocument') as doc_stream:
                    doc_content = doc_stream.read()
                    content = doc_content[:max_length].decode('utf-8', errors='ignore')
            else:
                raise ValueError("Unsupported or corrupted .doc file")

        elif file_extension == '.ppt':
            if olefile.isOleFile(file_path):
                ole = olefile.OleFileIO(file_path)
                with ole.openstream('PowerPoint Document') as ppt_stream:
                    ppt_content = ppt_stream.read()
                    content = ppt_content[:max_length].decode('utf-8', errors='ignore')
            else:
                raise ValueError("Unsupported or corrupted .ppt file")

        else:
            raise ValueError("Unsupported file type")

        content = content.replace(" ", "")

    except Exception as e:
        return f"An error occurred: {e}"

    kw = ""
    for index, item in enumerate(keyword_list, start=1):
        kw += f"{index}. {item} "
    kw.strip()

    llm = SparkLLM(appid, api_key, api_secret, Spark_url, domain)
    Input = input("\n" + "我:这里有一段文本内容：“" + content + "”，请根据其内容，从以下一组词汇中为其选定一个合适的主题:“" +
                  kw + "”，请给我一个数字作为回答。（注意，你的回答仅仅允许有一个数字！不要有其他任何多余的内容")

    ans = llm.query(Input)
    match = re.search(r'\d+', ans)
    if match:
        return int(match.group())
    else:
        return f"An error occurred: spark fail"
    